"""
Direct document generator — single LLM call that produces full, prompt-specific content.

Replaces the 7-agent pipeline which suffered from:
  • 7× LLM round-trips → rate-limit failures
  • Data-agent generating generic data unrelated to the user's request
  • Fallback structures overwriting real content when any one agent failed
  • "error.xlsx" / "generated_document.docx" fallback filenames

This module asks ONE model to produce the complete document spec (structure + content +
data) in one shot, then compiles it with lightweight inline compilers.
"""
from __future__ import annotations

import io
import json
import logging
import re
import textwrap
from typing import Any, Dict, List, Optional

from langchain_core.messages import HumanMessage, SystemMessage

from app.core.openrouter_client import openrouter_client

logger = logging.getLogger(__name__)


# ─────────────────────────────────────── prompts ──────────────────────────────

DETECT_FORMAT_PROMPT = """Determine the best document format for this request.
Reply with ONLY one word: excel, word, or powerpoint.

Request: {prompt}"""

EXCEL_PROMPT = """You are a professional data analyst and Excel expert.
Create a complete, realistic Excel spreadsheet for the following request.

REQUEST: {prompt}

Respond with VALID JSON ONLY — no markdown, no explanation, no code fences.
Use this exact schema:
{{
  "type": "excel",
  "filename": "Descriptive_Name.xlsx",
  "sheets": [
    {{
      "name": "Sheet Name",
      "headers": ["Col1", "Col2", "Col3"],
      "rows": [
        ["value", "value", "value"],
        ...
      ],
      "formulas": {{
        "D11": "=SUM(D2:D10)",
        "E11": "=AVERAGE(E2:E10)"
      }}
    }}
  ]
}}

Rules:
- filename must reflect the actual topic (e.g. "Q4_Sales_Report.xlsx", "Monthly_Budget.xlsx")
- Generate 12-20 realistic data rows that DIRECTLY answer the request — not generic sample data
- Column headers must match what the user asked for
- Include SUM/AVERAGE formulas in a totals row where appropriate
- Create multiple sheets if the topic needs it (e.g. Raw Data + Summary)
- Numbers should be realistic (e.g. sales figures in $1k-$50k range, percentages as 0-100)
"""

WORD_PROMPT = """You are a professional business writer.
Write a complete Word document for the following request.

REQUEST: {prompt}

Respond with VALID JSON ONLY — no markdown, no explanation, no code fences.
Use this exact schema:
{{
  "type": "word",
  "filename": "Descriptive_Name.docx",
  "title": "Document Title",
  "subtitle": "Optional subtitle",
  "sections": [
    {{
      "heading": "Section Heading",
      "level": 1,
      "content": "Full paragraph text. Write 3-5 sentences of substantive content relevant to the request. Do NOT use placeholder text.",
      "bullet_points": ["Specific point 1", "Specific point 2"]
    }}
  ],
  "tables": [
    {{
      "headers": ["Column1", "Column2", "Column3"],
      "rows": [["value", "value", "value"]]
    }}
  ]
}}

Rules:
- filename must reflect the actual topic (e.g. "Business_Proposal.docx", "Annual_Report_2024.docx")
- Write REAL, specific content in every "content" field — no lorem ipsum, no placeholders
- Create 5-8 sections with substantive paragraphs
- Include bullet_points where a list makes sense
- Add a table if data comparison is relevant
- The document must directly address the user's request
"""

POWERPOINT_PROMPT = """You are a professional presentation designer.
Create a complete PowerPoint presentation for the following request.

REQUEST: {prompt}

Respond with VALID JSON ONLY — no markdown, no explanation, no code fences.
Use this exact schema:
{{
  "type": "powerpoint",
  "filename": "Descriptive_Name.pptx",
  "title": "Presentation Title",
  "slides": [
    {{
      "type": "title",
      "title": "Presentation Title",
      "subtitle": "Subtitle or date"
    }},
    {{
      "type": "content",
      "title": "Slide Title",
      "bullets": ["Specific point 1", "Specific point 2", "Specific point 3"]
    }}
  ]
}}

Rules:
- filename must reflect the actual topic (e.g. "Product_Launch_Deck.pptx")
- Create 7-10 slides with specific, relevant content
- First slide must be type "title" with a descriptive subtitle
- Each content slide should have 3-5 bullet points with REAL information
- Include an agenda slide, key content slides, and a conclusion/CTA slide
- Bullets must be concise (max 12 words each) but informative
"""


# ─────────────────────────────────────── helpers ──────────────────────────────

def _extract_json_object(text: str) -> str:
    """Strip markdown fences and extract the outermost JSON object/array."""
    text = text.strip()
    # Remove code fences
    text = re.sub(r"^```(?:json)?", "", text, flags=re.MULTILINE).strip()
    text = re.sub(r"```$", "", text, flags=re.MULTILINE).strip()
    # Find first { to last }
    start = text.find("{")
    end = text.rfind("}")
    if start != -1 and end != -1:
        return text[start : end + 1]
    return text


def _slugify(text: str, max_len: int = 50) -> str:
    """Create a safe filename fragment from arbitrary text."""
    slug = re.sub(r"[^\w\s-]", "", text)
    slug = re.sub(r"[\s]+", "_", slug.strip())
    return slug[:max_len]


# ─────────────────────────────────── compilers ────────────────────────────────

def _compile_excel(spec: Dict[str, Any]) -> Dict[str, Any]:
    """Compile an Excel workbook from the spec dict."""
    from openpyxl import Workbook
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    from openpyxl.utils import get_column_letter

    wb = Workbook()
    sheets: List[Dict] = spec.get("sheets", [])
    if not sheets:
        sheets = [{"name": "Data", "headers": [], "rows": [], "formulas": {}}]

    thin = Side(style="thin", color="D0D0D0")
    border = Border(left=thin, right=thin, top=thin, bottom=thin)

    for idx, sheet_spec in enumerate(sheets):
        if idx == 0:
            ws = wb.active
            ws.title = str(sheet_spec.get("name", "Sheet1"))[:31]
        else:
            ws = wb.create_sheet(title=str(sheet_spec.get("name", f"Sheet{idx+1}"))[:31])

        headers: List[str] = sheet_spec.get("headers", [])
        rows: List[List] = sheet_spec.get("rows", [])
        formulas: Dict[str, str] = sheet_spec.get("formulas", {})

        # ── headers row ──
        for ci, h in enumerate(headers, 1):
            cell = ws.cell(row=1, column=ci, value=str(h))
            cell.font = Font(bold=True, color="FFFFFF", size=11, name="Calibri")
            cell.fill = PatternFill("solid", fgColor="2563EB")
            cell.alignment = Alignment(horizontal="center", vertical="center", wrap_text=True)
            cell.border = border

        # ── data rows ──
        for ri, row in enumerate(rows, 2):
            for ci, val in enumerate(row, 1):
                # Convert to proper Python type
                cell_val: Any = val
                if isinstance(val, str):
                    # Try numeric coercion
                    clean = val.replace(",", "").replace("$", "").replace("%", "").strip()
                    try:
                        cell_val = int(clean)
                    except ValueError:
                        try:
                            cell_val = float(clean)
                        except ValueError:
                            cell_val = val
                ws.cell(row=ri, column=ci, value=cell_val).border = border
                # Alternating row shading
                if ri % 2 == 0:
                    ws.cell(row=ri, column=ci).fill = PatternFill("solid", fgColor="F0F4FF")

        # ── formulas ──
        for cell_ref, formula in formulas.items():
            try:
                ws[cell_ref] = formula if formula.startswith("=") else f"={formula}"
                ws[cell_ref].font = Font(bold=True, name="Calibri")
                ws[cell_ref].fill = PatternFill("solid", fgColor="DBEAFE")
            except Exception as exc:
                logger.warning("Formula error %s=%s: %s", cell_ref, formula, exc)

        # ── auto-width ──
        for ci, header in enumerate(headers, 1):
            col_letter = get_column_letter(ci)
            max_w = len(str(header)) + 2
            for ri in range(2, min(len(rows) + 2, 22)):
                cv = ws.cell(row=ri, column=ci).value
                if cv is not None:
                    max_w = max(max_w, len(str(cv)) + 2)
            ws.column_dimensions[col_letter].width = min(max_w, 45)

        # Freeze header row
        ws.freeze_panes = "A2"

    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)
    raw = buf.getvalue()

    filename = spec.get("filename") or "spreadsheet.xlsx"
    if not filename.endswith(".xlsx"):
        filename += ".xlsx"

    return {
        "content": raw,
        "filename": filename,
        "mime_type": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "size": len(raw),
    }


def _compile_word(spec: Dict[str, Any]) -> Dict[str, Any]:
    """Compile a Word document from the spec dict."""
    from docx import Document
    from docx.shared import Pt, RGBColor, Inches
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    doc = Document()

    # ── page margins ──
    for section in doc.sections:
        section.top_margin = Inches(1)
        section.bottom_margin = Inches(1)
        section.left_margin = Inches(1.25)
        section.right_margin = Inches(1.25)

    # ── title ──
    title = spec.get("title", "Document")
    tp = doc.add_heading(title, level=0)
    tp.alignment = WD_ALIGN_PARAGRAPH.CENTER
    for run in tp.runs:
        run.font.color.rgb = RGBColor(0x1E, 0x40, 0xAF)  # blue-800

    subtitle = spec.get("subtitle", "")
    if subtitle:
        sp = doc.add_paragraph(subtitle)
        sp.alignment = WD_ALIGN_PARAGRAPH.CENTER
        sp.runs[0].font.size = Pt(13)
        sp.runs[0].font.italic = True
        sp.runs[0].font.color.rgb = RGBColor(0x64, 0x74, 0x8B)

    doc.add_paragraph()  # spacer

    # ── sections ──
    for sec in spec.get("sections", []):
        heading = sec.get("heading", "")
        level = min(int(sec.get("level", 1)), 3)
        content = sec.get("content", "")
        bullets = sec.get("bullet_points", []) or sec.get("bullets", [])

        if heading:
            h = doc.add_heading(heading, level=level)
            for run in h.runs:
                run.font.color.rgb = RGBColor(0x1E, 0x40, 0xAF)

        if content:
            p = doc.add_paragraph(content)
            p.paragraph_format.space_after = Pt(6)
            for run in p.runs:
                run.font.size = Pt(11)
                run.font.name = "Calibri"

        for bp in bullets:
            doc.add_paragraph(str(bp), style="List Bullet")

    # ── tables ──
    for tbl_spec in spec.get("tables", []):
        headers = tbl_spec.get("headers", [])
        rows = tbl_spec.get("rows", [])
        if not headers:
            continue
        tbl = doc.add_table(rows=1 + len(rows), cols=len(headers))
        tbl.style = "Table Grid"
        hdr_cells = tbl.rows[0].cells
        for ci, h in enumerate(headers):
            hdr_cells[ci].text = str(h)
            hdr_cells[ci].paragraphs[0].runs[0].font.bold = True
        for ri, row in enumerate(rows, 1):
            for ci, val in enumerate(row):
                if ci < len(tbl.rows[ri].cells):
                    tbl.rows[ri].cells[ci].text = str(val)
        doc.add_paragraph()

    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)
    raw = buf.getvalue()

    filename = spec.get("filename") or "document.docx"
    if not filename.endswith(".docx"):
        filename += ".docx"

    return {
        "content": raw,
        "filename": filename,
        "mime_type": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "size": len(raw),
    }


def _compile_powerpoint(spec: Dict[str, Any]) -> Dict[str, Any]:
    """Compile a PowerPoint presentation from the spec dict."""
    from pptx import Presentation as PPTXPresentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    prs = PPTXPresentation()
    TITLE_ONLY = 0        # layout index
    TITLE_CONTENT = 1

    slides_spec: List[Dict] = spec.get("slides", [])
    if not slides_spec:
        slides_spec = [
            {"type": "title", "title": spec.get("title", "Presentation"), "subtitle": ""},
        ]

    for slide_spec in slides_spec:
        stype = slide_spec.get("type", "content")

        if stype == "title":
            layout = prs.slide_layouts[TITLE_ONLY]
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = slide_spec.get("title", "")
            # Subtitle placeholder (index 1) may not exist in all templates
            try:
                slide.placeholders[1].text = slide_spec.get("subtitle", "")
            except (KeyError, IndexError):
                pass

        else:
            layout = prs.slide_layouts[TITLE_CONTENT]
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = slide_spec.get("title", "")

            bullets = slide_spec.get("bullets", []) or slide_spec.get("content", [])
            if isinstance(bullets, str):
                bullets = [bullets]

            try:
                tf = slide.placeholders[1].text_frame
                tf.clear()
                for i, bullet in enumerate(bullets):
                    if i == 0:
                        para = tf.paragraphs[0]
                    else:
                        para = tf.add_paragraph()
                    para.text = str(bullet)
                    para.level = 0
                    para.font.size = Pt(18)
                    para.font.name = "Calibri"
            except (KeyError, IndexError):
                pass

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    raw = buf.getvalue()

    filename = spec.get("filename") or "presentation.pptx"
    if not filename.endswith(".pptx"):
        filename += ".pptx"

    return {
        "content": raw,
        "filename": filename,
        "mime_type": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "size": len(raw),
    }


# ──────────────────────────────────── main API ────────────────────────────────

_FORMAT_PROMPTS = {
    "excel": EXCEL_PROMPT,
    "word": WORD_PROMPT,
    "powerpoint": POWERPOINT_PROMPT,
}

_COMPILERS = {
    "excel": _compile_excel,
    "word": _compile_word,
    "powerpoint": _compile_powerpoint,
}


async def detect_format(prompt: str) -> str:
    """Ask the LLM which document format fits best. Returns 'excel'|'word'|'powerpoint'."""
    # Fast keyword heuristics first (no LLM call needed)
    lower = prompt.lower()
    if any(w in lower for w in ["spreadsheet", "excel", "csv", "formula", "budget", "tracker",
                                 "pivot", "chart", "table", "data", "rows", "columns", "sheet"]):
        return "excel"
    if any(w in lower for w in ["presentation", "slides", "powerpoint", "pptx", "deck",
                                 "slideshow", "pitch"]):
        return "powerpoint"
    if any(w in lower for w in ["document", "report", "letter", "word", "docx", "proposal",
                                 "memo", "essay", "article", "write", "draft"]):
        return "word"

    # Fallback: quick LLM call
    try:
        msgs = [
            SystemMessage(content="You choose document formats. Reply with exactly one word."),
            HumanMessage(content=DETECT_FORMAT_PROMPT.format(prompt=prompt)),
        ]
        result = await openrouter_client.complete(msgs, model_preference="primary", temperature=0.1)
        result = result.strip().lower().split()[0]
        if result in _FORMAT_PROMPTS:
            return result
    except Exception:
        pass
    return "excel"  # safe default


async def generate(prompt: str) -> Dict[str, Any]:
    """
    Generate a document in a single LLM call and compile it.
    Returns the same dict as the existing compilers:
      { content, filename, mime_type, size }
    """
    fmt = await detect_format(prompt)
    generation_prompt = _FORMAT_PROMPTS[fmt].format(prompt=prompt)

    msgs = [
        SystemMessage(content="You are a professional document creator. Output valid JSON only."),
        HumanMessage(content=generation_prompt),
    ]

    last_exc: Optional[Exception] = None
    # Try each available model preference in order
    for pref in ("primary", "reasoning"):
        try:
            logger.info("DirectGenerator: requesting %s document (model_pref=%s)", fmt, pref)
            raw = await openrouter_client.complete(
                msgs,
                model_preference=pref,
                temperature=0.6,
                json_mode=True,
                max_tokens=4096,
            )
            json_str = _extract_json_object(raw)
            spec = json.loads(json_str)
            logger.info("DirectGenerator: JSON parsed OK, compiling %s", fmt)
            result = _COMPILERS[fmt](spec)
            logger.info("DirectGenerator: compiled → %s (%d bytes)", result["filename"], result["size"])
            return result
        except json.JSONDecodeError as exc:
            logger.warning("DirectGenerator JSON parse error (pref=%s): %s", pref, exc)
            last_exc = exc
        except Exception as exc:
            logger.error("DirectGenerator failed (pref=%s): %s", pref, exc, exc_info=True)
            last_exc = exc

    # Last-resort: generate a minimal document with the prompt as the filename hint
    logger.warning("DirectGenerator: all attempts failed, using minimal fallback")
    slug = _slugify(prompt)
    if fmt == "excel":
        spec = {
            "type": "excel",
            "filename": f"{slug}.xlsx",
            "sheets": [{"name": "Data", "headers": ["Item", "Value", "Notes"],
                        "rows": [["Sample 1", 100, ""],
                                 ["Sample 2", 200, ""],
                                 ["Sample 3", 300, ""]],
                        "formulas": {"B5": "=SUM(B2:B4)"}}],
        }
    elif fmt == "word":
        spec = {
            "type": "word",
            "filename": f"{slug}.docx",
            "title": prompt[:80],
            "sections": [{"heading": "Overview", "level": 1,
                          "content": f"This document was generated for: {prompt}",
                          "bullet_points": []}],
            "tables": [],
        }
    else:
        spec = {
            "type": "powerpoint",
            "filename": f"{slug}.pptx",
            "title": prompt[:80],
            "slides": [
                {"type": "title", "title": prompt[:80], "subtitle": "AI Generated"},
                {"type": "content", "title": "Overview",
                 "bullets": [f"Request: {prompt[:100]}",
                              "Please try again for more detailed content"]},
            ],
        }
    return _COMPILERS[fmt](spec)
