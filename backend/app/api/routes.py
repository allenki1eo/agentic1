"""FastAPI routes for the AI Office Suite."""
from fastapi import APIRouter, HTTPException, UploadFile, File, BackgroundTasks, Query
from fastapi.responses import StreamingResponse, Response
from typing import List, Optional
import json
import asyncio
from datetime import datetime
import uuid
import os
import io
import aiofiles

from app.models.schemas import (
    GenerationRequest, GenerationResponse, FormulaRequest, FormulaResponse,
    WorkflowState, StreamEvent, TaskType, OutputFormat
)
from app.agents.workflow import orchestrator
from app.agents.formula_agent import formula_agent
from app.agents.intent_classifier import intent_classifier
from app.agents.direct_generator import generate as direct_generate
from app.core.config import settings
import logging

logger = logging.getLogger(__name__)

router = APIRouter()

# In-memory storage for generated files and workflow states
# In production, use Redis or database
generated_files = {}
workflow_states = {}


@router.post("/generate", response_model=GenerationResponse)
async def generate_document(request: GenerationRequest):
    """Generate a document based on user request using direct single-call generation."""

    workflow_id = str(uuid.uuid4())[:8]

    try:
        # ── Primary path: single LLM call that generates full, prompt-specific content ──
        final_output = await direct_generate(request.prompt)

        if final_output.get("content"):
            file_id = str(uuid.uuid4())[:12]
            generated_files[file_id] = final_output

            return GenerationResponse(
                workflow_id=workflow_id,
                status="completed",
                message=f"Successfully generated {final_output.get('filename', 'document')}",
                download_url=f"/api/download/{file_id}",
                preview_data={
                    "filename": final_output.get("filename"),
                    "size": final_output.get("size"),
                    "mime_type": final_output.get("mime_type"),
                }
            )

        # ── Fallback: multi-agent pipeline ──
        logger.warning("DirectGenerator returned no content, falling back to orchestrator")
        workflow_state = await orchestrator.execute(
            user_input=request.prompt,
            file_ids=request.file_ids
        )
        workflow_states[workflow_state.workflow_id] = workflow_state

        if workflow_state.error:
            return GenerationResponse(
                workflow_id=workflow_state.workflow_id,
                status="failed",
                message=f"Generation failed: {workflow_state.error}",
                error=workflow_state.error
            )

        fallback_output = workflow_state.final_output or {}
        if fallback_output.get("content"):
            file_id = str(uuid.uuid4())[:12]
            generated_files[file_id] = fallback_output
            return GenerationResponse(
                workflow_id=workflow_state.workflow_id,
                status="completed",
                message=f"Successfully generated {fallback_output.get('filename', 'document')}",
                download_url=f"/api/download/{file_id}",
                preview_data={
                    "filename": fallback_output.get("filename"),
                    "size": fallback_output.get("size"),
                    "mime_type": fallback_output.get("mime_type"),
                }
            )

        return GenerationResponse(
            workflow_id=workflow_id,
            status="failed",
            message="Could not generate document. Please try rephrasing your request.",
            error="No output produced"
        )

    except Exception as e:
        import traceback
        logger.error(f"Generation failed: {traceback.format_exc()}")
        raise HTTPException(status_code=500, detail=str(e))


@router.post("/generate/stream")
async def generate_document_stream(request: GenerationRequest):
    """Generate a document with streaming progress updates."""
    
    async def event_generator():
        workflow_id = str(uuid.uuid4())[:8]
        progress_events: list = []

        # Send initial event
        yield f"data: {json.dumps({'type': 'start', 'workflow_id': workflow_id, 'timestamp': datetime.utcnow().isoformat()})}\n\n"

        try:
            # Progress callback — collects events; event_generator flushes them after execute()
            async def progress_callback(stage: str, data: dict):
                event = {
                    "type": "progress",
                    "workflow_id": workflow_id,
                    "stage": stage,
                    "data": data,
                    "timestamp": datetime.utcnow().isoformat()
                }
                progress_events.append(f"data: {json.dumps(event)}\n\n")

            # Execute workflow
            workflow_state = await orchestrator.execute(
                user_input=request.prompt,
                file_ids=request.file_ids,
                progress_callback=progress_callback
            )

            # Flush collected progress events
            for evt in progress_events:
                yield evt
            
            # Store and send completion
            workflow_states[workflow_id] = workflow_state
            
            final_output = workflow_state.final_output or {}
            if final_output.get("content"):
                file_id = str(uuid.uuid4())[:12]
                generated_files[file_id] = final_output
                
                completion_event = {
                    "type": "complete",
                    "workflow_id": workflow_id,
                    "status": "success",
                    "download_url": f"/api/download/{file_id}",
                    "filename": final_output.get("filename"),
                    "timestamp": datetime.utcnow().isoformat()
                }
            else:
                completion_event = {
                    "type": "complete",
                    "workflow_id": workflow_id,
                    "status": "error",
                    "error": "No output generated",
                    "timestamp": datetime.utcnow().isoformat()
                }
            
            yield f"data: {json.dumps(completion_event)}\n\n"
            
        except Exception as e:
            error_event = {
                "type": "error",
                "workflow_id": workflow_id,
                "error": str(e),
                "timestamp": datetime.utcnow().isoformat()
            }
            yield f"data: {json.dumps(error_event)}\n\n"
        
        # End stream
        yield "data: [DONE]\n\n"
    
    return StreamingResponse(
        event_generator(),
        media_type="text/event-stream",
        headers={
            "Cache-Control": "no-cache",
            "Connection": "keep-alive",
        }
    )


@router.post("/formula/generate", response_model=FormulaResponse)
async def generate_formula(request: FormulaRequest):
    """Generate an Excel formula from natural language description."""
    
    try:
        result = await formula_agent.generate_formula(request)
        return result
    except Exception as e:
        logger.error(f"Formula generation failed: {e}")
        raise HTTPException(status_code=500, detail=str(e))


@router.post("/classify")
async def classify_intent(request: dict):
    """Classify user intent from natural language."""
    
    try:
        user_input = request.get("prompt", "")
        file_names = request.get("file_names", [])
        
        intent = await intent_classifier.classify(user_input, file_names)
        return intent
    except Exception as e:
        logger.error(f"Intent classification failed: {e}")
        raise HTTPException(status_code=500, detail=str(e))


@router.get("/download/{file_id}")
async def download_file(file_id: str):
    """Download a generated file."""
    
    if file_id not in generated_files:
        raise HTTPException(status_code=404, detail="File not found")
    
    file_data = generated_files[file_id]
    
    return Response(
        content=file_data["content"],
        media_type=file_data.get("mime_type", "application/octet-stream"),
        headers={
            "Content-Disposition": f'attachment; filename="{file_data.get("filename", "download")}"'
        }
    )


@router.get("/workflow/{workflow_id}")
async def get_workflow_status(workflow_id: str):
    """Get status of a workflow."""
    
    if workflow_id not in workflow_states:
        raise HTTPException(status_code=404, detail="Workflow not found")
    
    return workflow_states[workflow_id]


@router.post("/upload")
async def upload_file(file: UploadFile = File(...)):
    """Upload a file for processing."""
    
    try:
        # Validate file
        file_ext = os.path.splitext(file.filename)[1].lower()
        if file_ext not in settings.ALLOWED_EXTENSIONS:
            raise HTTPException(
                status_code=400,
                detail=f"File type {file_ext} not allowed"
            )
        
        # Check file size
        content = await file.read()
        if len(content) > settings.MAX_FILE_SIZE_MB * 1024 * 1024:
            raise HTTPException(
                status_code=400,
                detail=f"File size exceeds {settings.MAX_FILE_SIZE_MB}MB limit"
            )
        
        # Generate file ID
        file_id = str(uuid.uuid4())[:12]
        
        # Store file (in production, use S3 or similar)
        upload_dir = "/tmp/uploads"
        os.makedirs(upload_dir, exist_ok=True)
        file_path = os.path.join(upload_dir, f"{file_id}{file_ext}")
        
        async with aiofiles.open(file_path, "wb") as f:
            await f.write(content)
        
        return {
            "file_id": file_id,
            "filename": file.filename,
            "size": len(content),
            "mime_type": file.content_type,
            "stored_path": file_path
        }
        
    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"File upload failed: {e}")
        raise HTTPException(status_code=500, detail=str(e))


@router.get("/preview/{file_id}")
async def preview_file(file_id: str, max_length: int = Query(5000, description="Maximum preview length")):
    """Extract and return text content preview from a generated file."""
    
    if file_id not in generated_files:
        raise HTTPException(status_code=404, detail="File not found")
    
    file_data = generated_files[file_id]
    mime_type = file_data.get("mime_type", "")
    content = file_data.get("content", b"")
    filename = file_data.get("filename", "")
    
    try:
        preview_data = {
            "filename": filename,
            "mime_type": mime_type,
            "size": len(content) if isinstance(content, bytes) else len(content.encode()),
            "content_preview": None,
            "structure_preview": None,
            "format": None
        }
        
        # Word document preview
        if "word" in mime_type or filename.endswith(".docx"):
            from docx import Document
            doc = Document(io.BytesIO(content) if isinstance(content, bytes) else io.BytesIO(content.encode()))
            
            full_text = []
            for para in doc.paragraphs:
                if para.text.strip():
                    full_text.append(para.text)
            
            preview_data["format"] = "word"
            preview_data["content_preview"] = "\n".join(full_text)[:max_length]
            
            # Extract document structure
            structure = {
                "title": full_text[0] if full_text else "Untitled",
                "paragraph_count": len(doc.paragraphs),
                "table_count": len(doc.tables),
                "sections": []
            }
            
            for para in doc.paragraphs[:20]:  # First 20 paragraphs
                if para.style and para.style.name and "Heading" in para.style.name:
                    structure["sections"].append({
                        "level": para.style.name,
                        "text": para.text[:100]
                    })
            
            preview_data["structure_preview"] = structure
            
        # Excel preview
        elif "spreadsheet" in mime_type or filename.endswith(".xlsx"):
            from openpyxl import load_workbook
            wb = load_workbook(io.BytesIO(content) if isinstance(content, bytes) else io.BytesIO(content.encode()))
            
            sheets_data = []
            for sheet_name in wb.sheetnames[:3]:  # First 3 sheets
                ws = wb[sheet_name]
                sheet_preview = {
                    "name": sheet_name,
                    "row_count": ws.max_row,
                    "column_count": ws.max_column,
                    "data": []
                }
                
                # First 10 rows
                for row in ws.iter_rows(min_row=1, max_row=10, values_only=True):
                    sheet_preview["data"].append([str(cell) if cell is not None else "" for cell in row])
                
                sheets_data.append(sheet_preview)
            
            preview_data["format"] = "excel"
            preview_data["structure_preview"] = {
                "sheet_count": len(wb.sheetnames),
                "sheets": sheets_data
            }
            preview_data["content_preview"] = f"Excel workbook with {len(wb.sheetnames)} sheets"
            
        # PowerPoint preview
        elif "presentation" in mime_type or filename.endswith(".pptx"):
            from pptx import Presentation
            prs = Presentation(io.BytesIO(content) if isinstance(content, bytes) else io.BytesIO(content.encode()))
            
            slides = []
            for i, slide in enumerate(prs.slides[:5]):  # First 5 slides
                slide_data = {
                    "slide_number": i + 1,
                    "layout": slide.slide_layout.name if slide.slide_layout else "Unknown",
                    "text": []
                }
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_data["text"].append(shape.text[:200])
                
                slides.append(slide_data)
            
            preview_data["format"] = "powerpoint"
            preview_data["structure_preview"] = {
                "slide_count": len(prs.slides),
                "slides": slides
            }
            preview_data["content_preview"] = f"Presentation with {len(prs.slides)} slides"
        
        else:
            preview_data["format"] = "unknown"
            preview_data["content_preview"] = "Preview not available for this file type"
        
        return preview_data
        
    except Exception as e:
        logger.error(f"Preview extraction failed: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to extract preview: {str(e)}")


@router.post("/edit/{file_id}", response_model=GenerationResponse)
async def edit_document(file_id: str, request: dict):
    """Edit an existing generated document using natural language instructions."""

    if file_id not in generated_files:
        raise HTTPException(status_code=404, detail="File not found")

    edit_instructions = request.get("instructions", "")
    if not edit_instructions:
        raise HTTPException(status_code=400, detail="No edit instructions provided")

    file_data = generated_files[file_id]
    original_filename = file_data.get("filename", "document")
    original_mime = file_data.get("mime_type", "")

    # Determine format hint from mime type
    if "spreadsheet" in original_mime or "excel" in original_mime:
        fmt_hint = "Excel spreadsheet"
        fmt_code = "excel"
    elif "wordprocessing" in original_mime or "word" in original_mime:
        fmt_hint = "Word document"
        fmt_code = "word"
    elif "presentation" in original_mime or "powerpoint" in original_mime:
        fmt_hint = "PowerPoint presentation"
        fmt_code = "powerpoint"
    else:
        fmt_hint = "document"
        fmt_code = "auto"

    # Build a new prompt that describes the edit
    edit_prompt = (
        f"Edit the existing {fmt_hint} named '{original_filename}'. "
        f"Apply the following modifications: {edit_instructions}. "
        f"Regenerate the complete updated {fmt_hint} with all changes applied."
    )

    try:
        workflow_state = await orchestrator.execute(
            user_input=edit_prompt,
            file_ids=[]
        )

        workflow_states[workflow_state.workflow_id] = workflow_state

        if workflow_state.error:
            return GenerationResponse(
                workflow_id=workflow_state.workflow_id,
                status="failed",
                message=f"Edit failed: {workflow_state.error}",
                error=workflow_state.error
            )

        final_output = workflow_state.final_output or {}
        if final_output.get("content"):
            new_file_id = str(uuid.uuid4())[:12]
            generated_files[new_file_id] = final_output

            return GenerationResponse(
                workflow_id=workflow_state.workflow_id,
                status="completed",
                message=f"Document updated: {final_output.get('filename', 'document')}",
                download_url=f"/api/download/{new_file_id}",
                preview_data={
                    "filename": final_output.get("filename"),
                    "size": final_output.get("size"),
                    "mime_type": final_output.get("mime_type")
                }
            )

        return GenerationResponse(
            workflow_id=workflow_state.workflow_id,
            status="completed",
            message="Edit completed but no output was produced"
        )

    except Exception as e:
        import traceback
        logger.error(f"Edit failed: {traceback.format_exc()}")
        raise HTTPException(status_code=500, detail=str(e))


@router.get("/test-generate")
async def test_generate():
    """Test the full document generation workflow."""
    try:
        workflow_state = await orchestrator.execute(
            user_input="Create a simple test document",
            file_ids=[]
        )
        
        return {
            "status": workflow_state.status,
            "workflow_id": workflow_state.workflow_id,
            "has_error": bool(workflow_state.error),
            "error": workflow_state.error,
            "has_final_output": bool(workflow_state.final_output),
            "final_output_keys": list(workflow_state.final_output.keys()) if workflow_state.final_output else [],
            "filename": workflow_state.final_output.get("filename") if workflow_state.final_output else None,
        }
    except Exception as e:
        import traceback
        return {
            "status": "error",
            "error": str(e),
            "traceback": traceback.format_exc()
        }


@router.get("/test-model")
async def test_model():
    """Test OpenRouter model connection."""
    from app.core.openrouter_client import openrouter_client
    from langchain_core.messages import HumanMessage, SystemMessage
    
    try:
        messages = [
            SystemMessage(content="You are a helpful assistant."),
            HumanMessage(content="Say 'OpenRouter connection successful' and nothing else.")
        ]
        
        response = await openrouter_client.complete(
            messages=messages,
            model_preference="primary",
            temperature=0.7
        )
        
        return {
            "status": "success",
            "model_used": openrouter_client.get_last_model(),
            "response": response,
            "timestamp": datetime.utcnow().isoformat()
        }
    except Exception as e:
        logger.error(f"Model test failed: {e}", exc_info=True)
        return {
            "status": "error",
            "error": str(e),
            "timestamp": datetime.utcnow().isoformat()
        }


@router.get("/health")
async def health_check():
    """Health check endpoint."""
    return {
        "status": "healthy",
        "timestamp": datetime.utcnow().isoformat(),
        "version": "1.0.0"
    }


@router.get("/models")
async def list_models():
    """List available AI models."""
    from app.core.openrouter_client import FREE_MODELS
    
    return {
        "primary": settings.PRIMARY_MODEL,
        "coding": settings.CODING_MODEL,
        "reasoning": settings.REASONING_MODEL,
        "fallbacks": settings.FALLBACK_MODELS,
        "all_free_models": [
            {"name": name, "context_length": config.context_length}
            for name, config in FREE_MODELS.items()
        ]
    }
